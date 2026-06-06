"""
Build Quest-On_Dongguk_Guide.pptx from rendered slide PNGs.

Each PPT slide is a single full-bleed PNG. No text/font dependencies in the
.pptx itself — fonts cannot break because PPT contains only raster images.
"""
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

HERE = Path(__file__).parent
SLIDES_DIR = HERE / "slides"
OUT = HERE / "Quest-On_Dongguk_Guide.pptx"

# 16:9 widescreen, matches deck (1280x720 → 13.333" x 7.5")
SLIDE_W = Emu(12192000)  # 13.333"
SLIDE_H = Emu(6858000)   # 7.5"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]  # blank layout

png_paths = sorted(SLIDES_DIR.glob("slide-*.png"))
if len(png_paths) != 10:
    raise SystemExit(f"Expected 10 slide PNGs, found {len(png_paths)}")

for png in png_paths:
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(
        str(png),
        left=Emu(0),
        top=Emu(0),
        width=SLIDE_W,
        height=SLIDE_H,
    )

prs.save(OUT)
size_mb = OUT.stat().st_size / 1_000_000
print(f"PPTX -> {OUT.name} ({size_mb:.2f} MB, {len(png_paths)} slides @ 13.333x7.5in)")
